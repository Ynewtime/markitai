#!/usr/bin/env python3
"""
Markdown 批量转换脚本

将文件或文件夹中的文档转换为 Markdown 格式，并提取图片到 assets 目录
支持子文件夹递归处理，保持原目录结构

支持格式: txt, doc, docx, ppt, pptx, xls, xlsx, pdf, png, jpg, html, csv, json, xml

使用方法:
    # 转换单个文件
    uv run convert_to_markdown.py "文件.docx"

    # 转换文件夹
    uv run convert_to_markdown.py "文件夹路径"

    # 启用 LLM 格式优化
    uv run convert_to_markdown.py -o "文件夹路径"

    # 分析文档中提取的附件图片
    uv run convert_to_markdown.py --analyze-attachments "文件夹路径"

    # 分析独立的图片文件
    uv run convert_to_markdown.py --analyze-image-files "文件夹路径"

    # 指定输出目录
    uv run convert_to_markdown.py -out "输出路径" "文件夹路径"

命令行参数:
    -o, --optimize            使用 LLM 优化 Markdown 格式
    --analyze-attachments     使用 LLM 分析文档中提取的附件图片
    --analyze-image-files     使用 LLM 分析独立的图片文件
    -out, --output            指定输出目录（默认在源文件同级或父级创建 output）
    -h, --help                显示帮助信息

说明:
    - 默认不使用 LLM 处理图片
    - 直接指定图片文件时，--analyze-image-files 默认启用
    - 所有选项可独立使用，也可组合使用

LLM 配置（可选，用于格式优化和图片分析）:
    设置环境变量:
    - OPENAI_API_KEY + OPENAI_MODEL: OpenAI
    - OPENAI_API_KEY + OPENAI_BASE_URL + OPENAI_MODEL: OpenRouter 等兼容接口
    - GOOGLE_API_KEY: Google Gemini
    - OLLAMA_MODEL: 本地 Ollama
    - LLM_PROMPT: 自定义图片描述提示词
    - LLM_CHUNK_SIZE: 超长文本跳过优化的阈值（默认 10000 字符）
    - MAX_WORKERS: 文件并行处理线程数（默认 10）
    - MAX_IMG_WORKERS: 单文件内图片并行分析线程数（默认 5）
    - OUTPUT_DIR: 默认输出目录

输出结构（文件夹模式）:
    输入: docs/my_data/
    输出: docs/output/      # 默认在父目录创建 output
    ├── 文件名.md
    ├── assets/
    │   └── 文件名_001.png
    └── 子文件夹/           # 保持原目录结构
        └── 文件名.md
"""

import base64
import concurrent.futures
import contextlib
import datetime
import io
import json
import logging
import os
import re
import shutil
import subprocess
import sys
import threading
import time
import uuid
from functools import lru_cache
from pathlib import Path

from markitdown import MarkItDown

# ============================================================
# 预编译正则表达式（性能优化）
# ============================================================
# Markdown 图片引用: ![alt](path)
RE_IMAGE_REF = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")
# Base64 内嵌图片: ![alt](data:image/xxx;base64,...)
RE_BASE64_IMAGE = re.compile(r"!\[([^\]]*)\]\(data:image/([^;]+);base64,([^)]+)\)")
# 文件名清理: 空格和特殊字符
RE_SANITIZE_CHARS = re.compile(r"[\s\(\)\[\]\{\}<>\'\"#%&!@\^\*\+\=\|\\:;,\?]+")
# 连续下划线
RE_MULTI_UNDERSCORE = re.compile(r"_+")
# 连续空白字符
RE_MULTI_WHITESPACE = re.compile(r"\s+")

# ============================================================
# 日志配置
# ============================================================
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S",
    handlers=[logging.StreamHandler(sys.stdout)],
)
logger = logging.getLogger("converter")

# 全局互斥锁
office_lock = threading.Lock()  # 用于保护 COM 对象
counter_lock = threading.Lock()  # 用于保护进度计数器
# 全局计数器
processed_count = 0
total_count = 0


def load_env_file():
    """从 .env 文件加载环境变量"""
    script_dir = Path(__file__).parent
    env_file = script_dir / ".env"

    if not env_file.exists():
        return

    with open(env_file, encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("#"):
                continue
            if "=" in line:
                key, _, value = line.partition("=")
                key = key.strip()
                value = value.strip().strip('"').strip("'")
                if key and value:
                    os.environ.setdefault(key, value)


# 加载 .env 文件
load_env_file()

# 支持的文件格式
SUPPORTED_FORMATS = {
    # 文档
    ".txt",
    ".doc",
    ".docx",
    ".ppt",
    ".pptx",
    ".xls",
    ".xlsx",
    ".pdf",
    # 图片（需要 LLM 才能生成描述）
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".webp",
    ".bmp",
    # 其他
    ".html",
    ".htm",
    ".csv",
    ".json",
    ".xml",
}
# 需要先转换的旧格式
OLD_FORMATS = {".doc": ".docx", ".ppt": ".pptx"}
# 图片格式
IMAGE_FORMATS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
# PPTX 格式（含嵌入图片，可选 LLM 分析）
PPTX_FORMATS = {".pptx"}

# LLM 内容描述提示词（用于图片/PPT/PDF等）
LLM_CONTENT_PROMPT = """请用中文对这个内容进行完整、结构化的描述，输出干净的 Markdown 格式。

## 输出要求

根据内容类型，按以下规范描述：

### 📊 如果是图表/数据可视化
- 图表类型（柱状图/折线图/饼图/流程图等）
- 标题和主题
- 坐标轴/图例说明
- 关键数据点和数值（尽可能提取具体数字）
- 数据趋势和结论

### 🖼️ 如果是普通图片/照片
- 主体内容和场景描述
- 重要的文字信息（完整提取）
- 颜色、构图、风格特点
- 图片可能的用途和上下文

### 📑 如果是文档/PPT页面
- 标题和章节结构
- 完整提取所有文字内容（不要遗漏）
- 要点列表和层级关系
- 表格数据（转为 Markdown 表格）
- 图表和图片的描述

### 📐 如果是技术图/架构图/流程图
- 图的类型和整体结构
- 各组件/节点的名称和含义
- 连接关系和数据流向
- 关键步骤和逻辑

## 格式规范

- 使用清晰的 Markdown 结构（标题、列表、表格）
- 标题从 ## 开始，不使用 #
- 数字和专有名词保持原样
- 提取的文字用 > 引用块标注
- 重要信息用 **加粗** 标注

## 核心原则

1. **完整性**：提取所有可见的文字和数据，不要省略
2. **准确性**：如实描述，不要推测不确定的内容
3. **结构化**：用合适的 Markdown 格式组织信息
4. **实用性**：描述应便于后续检索和理解"""

# 格式优化提示词（使用 {{PLACEHOLDER}} 占位符）
FORMAT_OPTIMIZE_PROMPT = """## 任务

将文件优化为高质量的 Markdown 格式。

## 输出要求

- 在头部添加 YAML Frontmatter，必须使用 `---` 包裹（严禁使用 ```yaml 代码块包裹）
- YAML Frontmatter 内容紧跟 `---` 后的下一行，不要有空行
- YAML Frontmatter 包含以下字段：
  - title: {{TITLE_INSTRUCTION}}
  - processed: {{PROCESSED_DATE}}
  - author: 作者（如果能识别，否则留空）
  - description: 文档内容摘要（如果包含特殊字符，请使用双引号包裹）
- 符合 Obsidian Flavored Markdown / GFM 规范

---

## 核心原则

仅做格式转换，严禁修改任何原文内容！不要总结、提炼或改写文案！
尽最大程度保留内容，对于长文本也不要省略内容，不要刻意删减内容

---

## 清洗规范 (Cleaning Rules)

1. **页眉页脚**
   - 删除重复出现的页眉页脚内容，如 `Page 1 of 10`
   - 删除公司机密声明，如 `Confidential`, `Internal Use Only` 等（如果它们作为页脚反复出现）
   - 删除无意义的分隔符或乱码

2. **空行与间距**
   - Header 与 Header 之间保留一个空行
   - Header 与正文之间保留一个空行
   - 连续的 <br/> 或空行合并为一个空行

## 格式规范

1. 标题
   - 不要自行添加标题层级
   - 标题层级从 h2 开始，不要使用 h1 (h1通常作为文件名/title)

2. 链接与图片
   - 链接：[text](href)
   - 图片：![alt](src)

3. 列表处理
   - 有序列表：保留数字序列，如 1. 项目
   - 无序列表：使用 - 或 *，如 - 项目
   - 嵌套列表：通过缩进表示层级

4. 文本样式
   - 加粗：**text**（处理规则见后文加粗文本部分）
   - 斜体：*text* 或 _text_
   - 删除线：~~text~~
   - 内联代码：`code`
   - 上标：^text^
   - 下标：~text~

5. 引用块
   - 使用 > 表示，多层引用使用 >>>
   - 示例：<blockquote>引用</blockquote> -> > 引用

6. 代码块
   - 标注正确的语言标识
   - 保持代码原样，不添加无关字符
   - 嵌套规则：外层代码块的反引号数量 = 内层最大反引号数 + 1

7. 表格
   - 使用 | 分隔列，- 分隔表头与内容
   - 示例：HTML 表格转换为 | 列1 | 列2 | 和 | --- | --- |

8. 分隔线
   - 使用 --- 或 *** 表示

9. 任务列表
   - 使用 - [ ] 和 - [x] 表示未完成和已完成任务

10. 脚注
    - 使用 [^1] 标注，并在文末定义 [^1]: 注释内容

11. 多媒体
    - 不支持 Markdown 的多媒体内容可保留 HTML 代码
    - 其他内容一律使用 Markdown 语法

12. 加粗文本
    - 连续的 <strong> 标签合并为单个 **text**
    - 标点位置修正：。** -> **。，，** -> **，，：** -> **：（重要！加粗的双星号要移到标点符号外）

---

## 待优化的 Markdown 内容

"""


def create_llm_client():
    """
    根据环境变量创建 LLM 客户端
    支持: OpenAI, OpenRouter, Google Gemini, Azure OpenAI, Ollama, 自定义接口

    优先级:
    1. OPENAI_API_BASE/OPENAI_BASE_URL (自定义接口，包括 OpenRouter)
    2. GOOGLE_API_KEY / GEMINI_API_KEY (Google Gemini)
    3. AZURE_OPENAI_* (Azure OpenAI)
    4. OPENAI_API_KEY (原生 OpenAI)
    5. OLLAMA_MODEL (本地 Ollama)
    """
    try:
        from openai import AzureOpenAI, OpenAI
    except ImportError:
        logger.warning("⚠️ 未安装 openai 包，LLM 功能不可用")
        logger.warning("   运行: uv add openai")
        return None, None

    # 1. 自定义 OpenAI 兼容接口（OpenRouter 等）
    api_base = os.environ.get("OPENAI_API_BASE") or os.environ.get("OPENAI_BASE_URL")
    if api_base:
        api_key = os.environ.get("OPENAI_API_KEY")
        model = os.environ.get("OPENAI_MODEL") or os.environ.get("LLM_MODEL", "gpt-4o")
        if api_key:
            client = OpenAI(api_key=api_key, base_url=api_base)
            provider = "OpenRouter" if "openrouter" in api_base.lower() else "自定义接口"
            logger.info(f"✅ 已启用 {provider} LLM: {model}")
            return client, model

    # 2. Google Gemini（通过 OpenAI 兼容接口）
    gemini_key = os.environ.get("GOOGLE_API_KEY") or os.environ.get("GEMINI_API_KEY")
    if gemini_key:
        model = os.environ.get("GEMINI_MODEL", "gemini-2.0-flash")
        client = OpenAI(
            api_key=gemini_key,
            base_url="https://generativelanguage.googleapis.com/v1beta/openai/",
        )
        logger.info(f"✅ 已启用 Google Gemini LLM: {model}")
        return client, model

    # 3. Azure OpenAI
    azure_key = os.environ.get("AZURE_OPENAI_API_KEY")
    azure_endpoint = os.environ.get("AZURE_OPENAI_ENDPOINT")
    if azure_key and azure_endpoint:
        client = AzureOpenAI(
            api_key=azure_key,
            api_version=os.environ.get("AZURE_OPENAI_API_VERSION", "2024-02-01"),
            azure_endpoint=azure_endpoint,
        )
        model = os.environ.get("AZURE_OPENAI_MODEL", "gpt-4o")
        logger.info(f"✅ 已启用 Azure OpenAI LLM: {model}")
        return client, model

    # 4. 原生 OpenAI
    openai_key = os.environ.get("OPENAI_API_KEY")
    if openai_key:
        model = os.environ.get("OPENAI_MODEL", "gpt-4o")
        client = OpenAI(api_key=openai_key)
        logger.info(f"✅ 已启用 OpenAI LLM: {model}")
        return client, model

    # 5. 本地 Ollama
    ollama_model = os.environ.get("OLLAMA_MODEL")
    if ollama_model:
        ollama_host = os.environ.get("OLLAMA_HOST", "http://localhost:11434")
        client = OpenAI(base_url=f"{ollama_host}/v1", api_key="ollama")
        logger.info(f"✅ 已启用 Ollama LLM: {ollama_model}")
        return client, ollama_model

    return None, None


def optimize_markdown_format(
    markdown_text: str,
    llm_client,
    llm_model: str,
    file_title: str | None = None,
    max_retries: int = 5,
    base_delay: float = 10.0,
    logger=print,
) -> str:
    """使用 LLM 优化 Markdown 格式，带重试机制

    Args:
        markdown_text: 待优化的 Markdown 文本
        llm_client: LLM 客户端
        llm_model: LLM 模型名称
        file_title: 文件标题（用于 Frontmatter）
        max_retries: 最大重试次数
        base_delay: 重试基础延迟（秒）
        logger: 日志函数

    注意：只发送纯文本给 LLM（包括 ![alt](path) 这样的图片引用），
    不会发送实际的图片数据。
    """
    if not llm_client or not llm_model:
        return markdown_text

    # 准备 Prompt，使用占位符替换
    current_prompt = FORMAT_OPTIMIZE_PROMPT

    # 替换 title 占位符
    if file_title:
        # 转义双引号，防止 YAML 语法错误
        safe_title = file_title.replace('"', '\\"')
        title_instruction = f'"{safe_title}" (请直接使用此标题)'
    else:
        title_instruction = '文件标题（从内容识别，如果包含特殊字符如冒号，请使用双引号包裹，例如 title: "Title: Subtitle"）'
    current_prompt = current_prompt.replace("{{TITLE_INSTRUCTION}}", title_instruction)

    # 准备日期占位符
    current_date = datetime.datetime.now().strftime("%Y-%m-%d")
    current_prompt = current_prompt.replace("{{PROCESSED_DATE}}", current_date)

    # 短文本：直接处理（保持原有逻辑，包含摘要生成）
    # 注意：日期已在函数开头替换到 current_prompt 中
    for attempt in range(max_retries):
        try:
            response = llm_client.chat.completions.create(
                model=llm_model,
                messages=[{"role": "user", "content": current_prompt + markdown_text}],
                temperature=0.1,  # 低温度，保持一致性
            )

            optimized = response.choices[0].message.content

            # 清理可能的 markdown 代码块包裹
            if optimized.startswith("```markdown"):
                optimized = optimized[len("```markdown") :].strip()
            elif optimized.startswith("```md"):
                optimized = optimized[len("```md") :].strip()
            # 注意：不建议直接清理 ```，因为可能误伤正文中的代码块
            # 但如果整段都是代码块包裹的，则需要清理
            elif optimized.startswith("```") and optimized.endswith("```"):
                optimized = optimized[3:-3].strip()

            # 修复：移除 Frontmatter 开始处多余的空行
            # 针对用户反馈: ---\n\ntitle: (--- 和 title 之间有空行)
            if optimized.startswith("---"):
                # 方法1: 移除 --- 后的所有空白行，直到遇到非空行
                lines = optimized.split("\n")
                if lines[0] == "---":
                    # 找到第一个非空行的索引
                    first_content_idx = 1
                    while first_content_idx < len(lines) and not lines[first_content_idx].strip():
                        first_content_idx += 1
                    # 重组: --- + 非空内容
                    if first_content_idx > 1:
                        optimized = "---\n" + "\n".join(lines[first_content_idx:])

            # 额外检查：如果头部是 ```yaml 包裹的 Frontmatter，尝试修复
            # 错误示例:
            # ```yaml
            # title: ...
            # ```
            # 正文...
            if optimized.startswith("```yaml"):
                optimized = optimized.replace("```yaml", "---", 1)
                # 寻找下一个 ``` 并替换为 ---
                if "\n```\n" in optimized:
                    optimized = optimized.replace("\n```\n", "\n---\n", 1)
                elif "\n```" in optimized:  # 处理紧凑情况
                    optimized = optimized.replace("\n```", "\n---", 1)

            # 额外检查：YAML 冒号转义修复
            # 有些模型生成 title: Something: Subtitle 而不是 title: "Something: Subtitle"
            try:
                # 提取 Frontmatter
                if optimized.startswith("---"):
                    end_idx = optimized.find("\n---", 3)
                    if end_idx != -1:
                        frontmatter = optimized[3:end_idx]
                        new_frontmatter = []
                        for line in frontmatter.splitlines():
                            # 跳过空行，防止在 --- 和 title 之间产生空行
                            if not line.strip():
                                continue
                            # 简单修复逻辑：检查常见字段，如果值包含冒号且未加引号，则添加引号
                            if ":" in line:
                                key, _, val = line.partition(":")
                                key = key.strip()
                                val = val.strip()
                                # 如果值包含冒号，且没有被引号包裹
                                if (
                                    ":" in val
                                    and not (val.startswith('"') and val.endswith('"'))
                                    and not (val.startswith("'") and val.endswith("'"))
                                ):
                                    # 针对 title 和 description 字段
                                    if key in ["title", "description"]:
                                        # 替换双引号为单引号，防止转义问题
                                        val_escaped = val.replace('"', '\\"')
                                        new_line = f'{key}: "{val_escaped}"'
                                        new_frontmatter.append(new_line)
                                        continue
                            new_frontmatter.append(line)

                        # 重组内容
                        optimized = "---\n" + "\n".join(new_frontmatter) + optimized[end_idx:]
            except Exception:
                pass  # 如果解析失败，保持原样

            return optimized

        except Exception as e:
            error_msg = str(e)

            # 检查是否是速率限制错误
            is_rate_limit = any(
                x in error_msg for x in ["429", "rate", "quota", "RESOURCE_EXHAUSTED", "Too Many"]
            )

            if is_rate_limit and attempt < max_retries - 1:
                # 指数退避：10s, 20s, 40s, 80s, 160s
                delay = base_delay * (2**attempt)
                logger(
                    f"      ⏳ API 速率限制，{delay:.0f}秒后重试 ({attempt + 1}/{max_retries})..."
                )
                time.sleep(delay)
            else:
                logger(f"      ⚠️ 格式优化失败: {e}")
                return markdown_text

    return markdown_text


# 图片分析提示词
IMAGE_ANALYZE_PROMPT = """请分析这张图片的内容，并以 JSON 格式返回两份描述。

要求返回 JSON 格式如下：
{
    "summary": "一句话总结图片内容，简明扼要，不换行，不使用 Markdown 语法，用于 alt 属性。",
    "detail": "详细描述图片内容，覆盖所有细节（文字、数据、颜色、布局等）。使用 Markdown 格式（可以包含标题、列表、表格等），用于单独的文档。"
}

注意：
1. 仅返回合法的 JSON 字符串，不要包含 ```json 包裹。
2. "summary" 字段必须简洁，适合作为图片的替代文本。
3. "detail" 字段要尽可能详细，不要遗漏信息。
"""


def analyze_image_with_llm(
    image_path: Path,
    llm_client,
    llm_model: str,
    max_retries: int = 5,
    base_delay: float = 10.0,
    logger=print,
) -> str | None:
    """使用 LLM 分析图片，生成描述文本

    Returns:
        str: 净化后的 summary（用于 alt 文本）
    Side Effect:
        在 assets 目录下生成同名的 .md 文件，包含 detail 详细描述
    """
    if not llm_client or not llm_model or not image_path.exists():
        return None

    # base64, json, time 已在文件顶部导入

    # 读取图片并转为 base64
    try:
        image_data = image_path.read_bytes()
        base64_image = base64.b64encode(image_data).decode("utf-8")

        # 根据扩展名确定 MIME 类型
        ext = image_path.suffix.lower()
        mime_map = {
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".png": "image/png",
            ".gif": "image/gif",
            ".webp": "image/webp",
            ".bmp": "image/bmp",
        }
        mime_type = mime_map.get(ext, "image/png")

    except Exception as e:
        logger(f"        ⚠️ 读取图片失败: {e}")
        return None

    for attempt in range(max_retries):
        try:
            # 准备请求参数
            request_params = {
                "model": llm_model,
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": IMAGE_ANALYZE_PROMPT},
                            {
                                "type": "image_url",
                                "image_url": {"url": f"data:{mime_type};base64,{base64_image}"},
                            },
                        ],
                    }
                ],
                # 移除 max_tokens 限制，避免长描述截断导致 JSON 格式错误
                # "max_tokens": 2000,
            }

            # 尝试添加 response_format 参数（OpenAI/Azure/Ollama JSON 模式）
            # 注意：某些模型/旧版 OpenAI 接口可能不支持此参数，如果报错会回退
            try:
                request_params["response_format"] = {"type": "json_object"}
                response = llm_client.chat.completions.create(**request_params)
            except Exception as e:
                # 如果是参数错误（400 Bad Request），尝试移除 response_format 重试
                error_str = str(e).lower()
                if (
                    "response_format" in error_str
                    or "unsupported parameter" in error_str
                    or "400" in error_str
                ):
                    # logger(f"        ℹ️ 模型不支持 response_format，尝试普通文本模式")
                    del request_params["response_format"]
                    response = llm_client.chat.completions.create(**request_params)
                else:
                    raise e  # 其他错误（如 429）抛出给外层重试逻辑处理

            content = response.choices[0].message.content.strip()

            # 清理可能的 Markdown 代码块包裹
            if content.startswith("```json"):
                content = content[7:]
            elif content.startswith("```"):
                content = content[3:]
            if content.endswith("```"):
                content = content[:-3]

            content = content.strip()

            try:
                data = json.loads(content)
                summary = data.get("summary", "")
                detail = data.get("detail", "")
            except json.JSONDecodeError:
                # 降级处理：如果不是 JSON，假设整个内容是 detailed，尝试生成一个 summary
                logger("        ⚠️ 响应非 JSON 格式，尝试降级处理")
                detail = content
                summary = content[:50].replace("\n", " ") + "..."

            # 1. 保存详细描述到 .md 文件
            if detail:
                try:
                    # 使用 <图片名>.<图片后缀>.md 格式
                    desc_file = image_path.with_name(f"{image_path.name}.md")
                    desc_file.write_text(detail, encoding="utf-8")
                except Exception as e:
                    logger(f"        ⚠️ 保存描述文件失败: {e}")

            # 2. 返回 Summary 用于 alt
            if summary:
                # 再次确保 summary 无换行（使用预编译正则）
                clean_summary = summary.replace("\n", " ").replace("\r", " ")
                clean_summary = RE_MULTI_WHITESPACE.sub(" ", clean_summary).strip()
                return clean_summary

            return None

        except Exception as e:
            error_msg = str(e)

            is_rate_limit = any(
                x in error_msg for x in ["429", "rate", "quota", "RESOURCE_EXHAUSTED", "Too Many"]
            )

            if is_rate_limit and attempt < max_retries - 1:
                delay = base_delay * (2**attempt)
                logger(
                    f"        ⏳ API 速率限制，{delay:.0f}秒后重试 ({attempt + 1}/{max_retries})..."
                )
                time.sleep(delay)
            else:
                logger(f"        ⚠️ 图片分析失败: {e}")
                return None

    return None


def analyze_images_in_markdown(
    markdown_text: str,
    assets_dir: Path,
    llm_client,
    llm_model: str,
    use_concurrency: bool = True,
    logger=logger.info,
) -> str:
    """分析 Markdown 中的图片，使用线程池并行处理

    Args:
        use_concurrency: 是否使用多线程并行（当外层已经并行时，内层建议串行以防 API 洪泛）
        logger: 日志函数
    """
    if not llm_client or not llm_model:
        return markdown_text

    # 匹配 ![alt](path) 格式（使用预编译正则）
    matches = list(RE_IMAGE_REF.finditer(markdown_text))

    if not matches:
        return markdown_text

    logger(f"      🔍 发现 {len(matches)} 张图片，正在分析...")

    # 准备任务
    tasks = []
    for match in matches:
        img_path_str = match.group(2)

        # 处理相对路径
        if img_path_str.startswith("assets/"):
            img_path = assets_dir.parent / img_path_str
        else:
            img_path = Path(img_path_str)

        tasks.append((match, img_path, img_path_str))

    replacements = []

    # 决定是否并行
    if use_concurrency:
        # 从环境变量获取并发数，默认为 5
        max_workers = int(os.environ.get("MAX_IMG_WORKERS", "5"))
    else:
        max_workers = 1

    # 使用 ThreadPoolExecutor，即使 max_workers=1 也可以统一代码结构
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        # 提交任务
        future_to_match = {
            executor.submit(
                analyze_image_with_llm,
                task[1],  # img_path
                llm_client,
                llm_model,
                5,  # max_retries
                2.0,  # base_delay (并发时稍微增加退避基数)
                logger,
            ): task
            for task in tasks
        }

        # 获取结果
        for completed_count, future in enumerate(
            concurrent.futures.as_completed(future_to_match), 1
        ):
            match, img_path, img_path_str = future_to_match[future]

            try:
                description = future.result()
                if description:
                    # 记录替换信息：(start_index, end_index, new_text)
                    new_ref = f"![{description}]({img_path_str})"
                    replacements.append((match.start(), match.end(), new_ref))
                    logger(
                        f"        ✅ [{completed_count}/{len(matches)}] {description[:30]}... ({img_path.name})"
                    )
                else:
                    logger(
                        f"        ⏭️  [{completed_count}/{len(matches)}] 无描述 ({img_path.name})"
                    )
            except Exception as e:
                logger(f"        ⚠️  [{completed_count}/{len(matches)}] 分析异常: {e}")

    # 按位置从后往前替换，避免索引偏移
    replacements.sort(key=lambda x: x[0], reverse=True)

    result = markdown_text
    for start, end, new_text in replacements:
        result = result[:start] + new_text + result[end:]

    return result


@lru_cache(maxsize=1)
def find_libreoffice() -> str | None:
    """查找 LibreOffice 可执行文件路径"""
    possible_paths = [
        # Windows
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        # Linux
        "/usr/bin/soffice",
        "/usr/bin/libreoffice",
        # macOS
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]

    for path in possible_paths:
        if os.path.exists(path):
            return path

    for cmd in ["soffice", "libreoffice"]:
        if shutil.which(cmd):
            return cmd

    return None


def check_ms_office_available() -> bool:
    """检查 Windows 上是否安装了 MS Office PowerPoint"""
    if sys.platform != "win32":
        return False

    try:
        # 通过注册表检查 PowerPoint 是否安装
        import winreg

        try:
            key = winreg.OpenKey(winreg.HKEY_CLASSES_ROOT, r"PowerPoint.Application")
            winreg.CloseKey(key)
            return True
        except OSError:  # WindowsError 是 OSError 的别名，使用 OSError 更通用
            return False
    except ImportError:
        return False


def check_ms_word_available() -> bool:
    """检查 Windows 上是否安装了 MS Office Word"""
    if sys.platform != "win32":
        return False

    try:
        import winreg

        try:
            key = winreg.OpenKey(winreg.HKEY_CLASSES_ROOT, r"Word.Application")
            winreg.CloseKey(key)
            return True
        except OSError:  # WindowsError 是 OSError 的别名，使用 OSError 更通用
            return False
    except ImportError:
        return False


def convert_with_ms_office(input_file: Path, output_dir: Path) -> Path | None:
    """使用 MS Office PowerPoint 转换 .ppt 到 .pptx（仅 Windows）"""
    if sys.platform != "win32":
        return None

    output_file = output_dir / (input_file.stem + ".pptx")

    # 转义路径中的单引号，防止 PowerShell 语法错误
    input_path = str(input_file.resolve()).replace("'", "''")
    output_path = str(output_file.resolve()).replace("'", "''")

    # 使用 PowerShell 调用 COM 对象，避免额外依赖
    ps_script = f"""
$ppt = New-Object -ComObject PowerPoint.Application
$ppt.Visible = [Microsoft.Office.Core.MsoTriState]::msoFalse
try {{
    $presentation = $ppt.Presentations.Open('{input_path}', $true, $false, $false)
    $presentation.SaveAs('{output_path}', 24)  # 24 = ppSaveAsOpenXMLPresentation
    $presentation.Close()
    Write-Host "SUCCESS"
}} catch {{
    Write-Host "FAILED: $_"
}} finally {{
    $ppt.Quit()
    [System.Runtime.Interopservices.Marshal]::ReleaseComObject($ppt) | Out-Null
}}
"""

    try:
        result = subprocess.run(
            ["powershell", "-NoProfile", "-Command", ps_script],
            capture_output=True,
            text=True,
            timeout=120,
        )

        if "SUCCESS" in result.stdout and output_file.exists():
            return output_file
        else:
            return None

    except Exception:
        return None


def convert_doc_with_ms_word(input_file: Path, output_dir: Path) -> Path | None:
    """使用 MS Office Word 转换 .doc 到 .docx（仅 Windows）"""
    if sys.platform != "win32":
        return None

    output_file = output_dir / (input_file.stem + ".docx")

    # 转义路径中的单引号，防止 PowerShell 语法错误
    input_path = str(input_file.resolve()).replace("'", "''")
    output_path = str(output_file.resolve()).replace("'", "''")

    # 使用 PowerShell 调用 COM 对象
    ps_script = f"""
$word = New-Object -ComObject Word.Application
$word.Visible = $false
try {{
    $doc = $word.Documents.Open('{input_path}')
    $doc.SaveAs2('{output_path}', 16)  # 16 = wdFormatDocumentDefault (.docx)
    $doc.Close()
    Write-Host "SUCCESS"
}} catch {{
    Write-Host "FAILED: $_"
}} finally {{
    $word.Quit()
    [System.Runtime.Interopservices.Marshal]::ReleaseComObject($word) | Out-Null
}}
"""

    try:
        result = subprocess.run(
            ["powershell", "-NoProfile", "-Command", ps_script],
            capture_output=True,
            text=True,
            timeout=120,
        )

        if "SUCCESS" in result.stdout and output_file.exists():
            return output_file
        else:
            return None

    except Exception:
        return None


def convert_with_libreoffice(input_file: Path, output_dir: Path, new_ext: str) -> Path | None:
    """使用 LibreOffice 转换旧格式"""
    libreoffice = find_libreoffice()
    if not libreoffice:
        return None

    cmd = [
        libreoffice,
        "--headless",
        "--convert-to",
        new_ext[1:],  # 去掉点号
        "--outdir",
        str(output_dir),
        str(input_file),
    ]

    try:
        subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        new_file = output_dir / (input_file.stem + new_ext)
        if new_file.exists():
            return new_file
        return None
    except Exception:
        return None


def convert_old_format(input_file: Path, temp_dir: Path, logger=logger.info) -> Path | None:
    """转换旧格式文件（.ppt/.doc 等）"""
    suffix = input_file.suffix.lower()
    new_ext = OLD_FORMATS.get(suffix)

    if not new_ext:
        return input_file

    temp_dir.mkdir(parents=True, exist_ok=True)

    # 对于 .ppt 文件，优先使用 MS Office（Windows）
    if suffix == ".ppt":
        # 1. 尝试 MS Office（Windows 优先）
        if check_ms_office_available():
            logger("  🔄 使用 MS Office 转换 .ppt → .pptx ...")
            result = convert_with_ms_office(input_file, temp_dir)
            if result:
                return result
            logger("      ⚠️ MS Office 转换失败，尝试 LibreOffice...")

        # 2. 尝试 LibreOffice
        libreoffice = find_libreoffice()
        if libreoffice:
            logger("  🔄 使用 LibreOffice 转换 .ppt → .pptx ...")
            result = convert_with_libreoffice(input_file, temp_dir, new_ext)
            if result:
                return result
            logger("      ⚠️ LibreOffice 转换失败")

        # 3. 都失败，提示用户手动转换
        logger("  ⚠️ 无法自动转换 .ppt 文件")
        logger("      📌 建议：用 PowerPoint 打开后另存为 .pptx 格式")
        logger("      （手动转换后可保留完整的图片和格式）")
        return None

    # 对于 .doc 文件
    if suffix == ".doc":
        # 1. 尝试 MS Word（Windows 优先）
        if check_ms_word_available():
            logger("  🔄 使用 MS Word 转换 .doc → .docx ...")
            result = convert_doc_with_ms_word(input_file, temp_dir)
            if result:
                return result
            logger("      ⚠️ MS Word 转换失败，尝试 LibreOffice...")

        # 2. 尝试 LibreOffice
        libreoffice = find_libreoffice()
        if libreoffice:
            logger("  🔄 使用 LibreOffice 转换 .doc → .docx ...")
            result = convert_with_libreoffice(input_file, temp_dir, new_ext)
            if result:
                return result
            logger("      ⚠️ LibreOffice 转换失败")

        # 3. 都失败，提示用户手动转换
        logger("  ⚠️ 无法自动转换 .doc 文件")
        logger("      📌 建议：用 Word 打开后另存为 .docx 格式")
        logger("      （手动转换后可保留完整的图片和格式）")
        return None

    return None


def sanitize_filename(name: str) -> str:
    """将文件名中的空格和特殊字符替换为下划线，确保 Markdown 兼容"""
    # 替换空格和常见特殊字符为下划线（使用预编译正则）
    result = RE_SANITIZE_CHARS.sub("_", name)
    # 合并连续的下划线（使用预编译正则）
    result = RE_MULTI_UNDERSCORE.sub("_", result)
    # 去除首尾下划线
    result = result.strip("_")
    return result if result else "image"


def extract_base64_images(
    markdown_text: str, assets_dir: Path, file_stem: str, logger=logger.info
) -> str:
    """从 Markdown 中提取 base64 图片，保存到 assets 目录，替换为相对路径

    对于 EMF/WMF 等 Markdown 不支持的格式，尝试转换为 PNG，失败则跳过
    """
    # 使用预编译的正则表达式

    # 清理文件名前缀
    safe_stem = sanitize_filename(file_stem)
    img_count = 0
    skip_count = 0

    # Markdown/浏览器支持的图片格式
    SUPPORTED_IMG_FORMATS = {
        "jpeg",
        "jpg",
        "png",
        "gif",
        "webp",
        "bmp",
        "svg+xml",
        "svg",
    }
    # 需要转换的格式
    CONVERTIBLE_FORMATS = {"x-emf", "emf", "x-wmf", "wmf", "tiff", "tif"}

    def try_convert_to_png(img_data: bytes, original_format: str) -> bytes | None:
        """尝试将不支持的图片格式转换为 PNG"""
        try:
            # io 已在文件顶部导入
            from PIL import Image

            # 尝试用 Pillow 打开并转换
            img = Image.open(io.BytesIO(img_data))

            # 转换为 RGB（处理 RGBA、P 等模式）
            if img.mode in ("RGBA", "LA", "P"):
                # 保持透明度
                img = img.convert("RGBA")
            elif img.mode != "RGB":
                img = img.convert("RGB")

            # 保存为 PNG
            output = io.BytesIO()
            img.save(output, format="PNG")
            return output.getvalue()

        except Exception as e:
            # Pillow 失败，如果是 EMF/WMF，尝试使用 LibreOffice 转换
            # LibreOffice 支持 headless 转换: soffice --headless --convert-to png --outdir ... file.emf
            is_emf = original_format.lower() in ("emf", "wmf", "x-emf", "x-wmf")
            libreoffice = find_libreoffice()

            if is_emf and libreoffice:
                try:
                    # 创建临时文件（uuid 已在文件顶部导入）
                    temp_name = f"temp_{uuid.uuid4().hex[:8]}"
                    # 去掉 x- 前缀，libreoffice 可能更喜欢标准扩展名
                    ext = original_format.replace("x-", "")
                    temp_emf = assets_dir / f"{temp_name}.{ext}"
                    temp_png = assets_dir / f"{temp_name}.png"

                    assets_dir.mkdir(parents=True, exist_ok=True)
                    temp_emf.write_bytes(img_data)

                    cmd = [
                        libreoffice,
                        "--headless",
                        "--convert-to",
                        "png",
                        "--outdir",
                        str(assets_dir),
                        str(temp_emf),
                    ]

                    # 转换
                    subprocess.run(cmd, capture_output=True, timeout=30)

                    if temp_png.exists():
                        png_data = temp_png.read_bytes()
                        # 清理
                        try:
                            temp_emf.unlink(missing_ok=True)
                            temp_png.unlink(missing_ok=True)
                        except Exception:
                            pass
                        return png_data

                    # 清理失败的临时文件
                    with contextlib.suppress(Exception):
                        temp_emf.unlink(missing_ok=True)

                except Exception as e2:
                    logger(f"        ⚠️ LibreOffice 转换失败: {e2}")

            logger(f"        ⚠️ 格式转换失败 ({original_format}): {e}")
            return None

    def replace_image(match):
        nonlocal img_count, skip_count

        alt_text = match.group(1)
        img_format = match.group(2).lower()
        base64_data = match.group(3)

        try:
            img_data = base64.b64decode(base64_data)
        except Exception as e:
            logger(f"      ⚠️ base64 解码失败: {e}")
            skip_count += 1
            return ""  # 移除无效图片

        # 检查是否需要转换
        if img_format in CONVERTIBLE_FORMATS:
            # logger(f"      🔄 转换 {img_format} 格式...")
            converted_data = try_convert_to_png(img_data, img_format)
            if converted_data:
                img_data = converted_data
                img_format = "png"
            else:
                logger(f"      ⏭️  跳过不支持的格式: {img_format}")
                skip_count += 1
                return ""  # 移除无法转换的图片

        # 检查是否是支持的格式
        if img_format not in SUPPORTED_IMG_FORMATS:
            logger(f"      ⏭️  跳过不支持的格式: {img_format}")
            skip_count += 1
            return ""  # 移除不支持的格式

        img_count += 1

        ext_map = {
            "jpeg": ".jpg",
            "jpg": ".jpg",
            "png": ".png",
            "gif": ".gif",
            "webp": ".webp",
            "bmp": ".bmp",
            "svg+xml": ".svg",
            "svg": ".svg",
        }
        ext = ext_map.get(img_format, f".{img_format}")

        img_filename = f"{safe_stem}_{img_count:03d}{ext}"
        img_path = assets_dir / img_filename

        try:
            assets_dir.mkdir(parents=True, exist_ok=True)
            img_path.write_bytes(img_data)
            # logger(f"      📷 {img_filename}")
        except Exception as e:
            logger(f"      ⚠️ 图片保存失败: {e}")
            return match.group(0)

        return f"![{alt_text}](assets/{img_filename})"

    result = RE_BASE64_IMAGE.sub(replace_image, markdown_text)
    return result


def extract_pptx_text_fallback(input_file: Path) -> str | None:
    """当 MarkItDown 失败时，使用 python-pptx 直接提取文本（不含图片）

    注意：python-pptx 已随 markitdown[all] 安装，无需额外安装
    """
    try:
        from pptx import Presentation

        prs = Presentation(str(input_file))
        markdown_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            markdown_parts.append(f"\n## 幻灯片 {slide_num}\n\n")

            for shape in slide.shapes:
                # 提取文本框内容
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:  # type: ignore
                        text = paragraph.text.strip()
                        if text:
                            # 尝试根据字体大小判断层级
                            try:
                                if paragraph.runs and paragraph.runs[0].font.size:
                                    font_size = paragraph.runs[0].font.size.pt
                                    if font_size and font_size >= 24:
                                        markdown_parts.append(f"### {text}\n\n")
                                    elif font_size and font_size >= 18:
                                        markdown_parts.append(f"**{text}**\n\n")
                                    else:
                                        markdown_parts.append(f"- {text}\n")
                                else:
                                    markdown_parts.append(f"- {text}\n")
                            except Exception:
                                markdown_parts.append(f"- {text}\n")

                # 提取表格
                if shape.has_table:
                    table = shape.table  # type: ignore
                    markdown_parts.append("\n")
                    for row_idx, row in enumerate(table.rows):
                        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        markdown_parts.append("| " + " | ".join(cells) + " |\n")
                        if row_idx == 0:
                            markdown_parts.append("| " + " | ".join(["---"] * len(cells)) + " |\n")
                    markdown_parts.append("\n")

            markdown_parts.append("\n")

        return "".join(markdown_parts)

    except Exception as e:
        logger.error(f"      ⚠️ 备选提取失败: {e}")
        return None


def convert_to_markdown(
    input_file: Path,
    output_dir: Path,
    assets_dir: Path,
    md: MarkItDown,
    llm_client=None,
    llm_model: str | None = None,
    optimize: bool | None = False,
    analyze_attachments: bool | None = False,
    output_filename: str | None = None,
    logger=logger.info,
) -> bool:
    """将文件转换为 Markdown，提取图片，可选格式优化和附件图片分析

    Args:
        md: MarkItDown 实例（根据文件类型可能带或不带 LLM 配置）
        analyze_attachments: 使用 LLM 分析文档中提取出的附件图片
        output_filename: 指定输出文件名（可选，用于防重名）
        logger: 日志处理函数
    """
    markdown_text = None

    # 第一次尝试：带图片提取
    try:
        result = md.convert(str(input_file), keep_data_uris=True)
        markdown_text = result.text_content

        # 提取 base64 图片
        markdown_text = extract_base64_images(
            markdown_text, assets_dir, input_file.stem, logger=logger
        )

        # 处理 Excel 转换结果中的 NaN 值（简单策略：将表格中的 NaN 替换为空白）
        if input_file.suffix.lower() in (".xls", ".xlsx") and markdown_text:
            markdown_text = re.sub(r"\|\s*NaN\s*\|", "||", markdown_text)

    except Exception as e:
        error_msg = str(e)

        # 检查是否是图片识别错误（常见于旧版 PPT 的 WMF/EMF 图片）
        if "UnidentifiedImageError" in error_msg or "cannot identify image file" in error_msg:
            logger("  ⚠️ 包含无法识别的图片格式，使用纯文本提取...")

            # 对于 PPTX 文件，使用 python-pptx 备选方案
            if input_file.suffix.lower() in (".pptx",):
                markdown_text = extract_pptx_text_fallback(input_file)
                if not markdown_text:
                    return False
            else:
                # 其他格式尝试不带图片提取
                try:
                    result = md.convert(str(input_file))
                    markdown_text = result.text_content
                except Exception as e2:
                    logger(f"  ❌ 转换失败: {e2}")
                    return False
        else:
            logger(f"  ❌ 转换失败: {e}")
            return False

    # 检查文本长度，如果超过阈值则跳过所有 LLM 处理（图片分析和格式优化）
    if (optimize or analyze_attachments) and llm_client and llm_model and markdown_text:
        DEFAULT_CHUNK_SIZE = 10000
        env_chunk_size = os.environ.get("LLM_CHUNK_SIZE")
        try:
            CHUNK_SIZE = int(env_chunk_size) if env_chunk_size else DEFAULT_CHUNK_SIZE
        except ValueError:
            CHUNK_SIZE = DEFAULT_CHUNK_SIZE

        text_length = len(markdown_text)
        if text_length > CHUNK_SIZE:
            logger(
                f"      ⚠️ 文本过长 ({text_length} > {CHUNK_SIZE})，跳过所有 LLM 处理（分析与优化）"
            )
            analyze_attachments = False
            optimize = False

    # 分析提取出的附件图片（在格式优化之前）
    if analyze_attachments and llm_client and llm_model and markdown_text:
        # 修正：用户要求宽松的优化，所以这里启用并发。
        # analyze_images_in_markdown 会读取 MAX_IMG_WORKERS (默认5)。
        markdown_text = analyze_images_in_markdown(
            markdown_text,
            assets_dir,
            llm_client,
            llm_model,
            use_concurrency=True,
            logger=logger,
        )

    # 格式优化（如果启用且配置了 LLM）
    if optimize and llm_client and llm_model and markdown_text:
        logger("      🔄 格式优化中...")
        markdown_text = optimize_markdown_format(
            markdown_text,
            llm_client,
            llm_model,
            file_title=input_file.stem,
            logger=logger,
        )

    # 写入 Markdown 文件
    if markdown_text is not None:
        if output_filename:
            output_file = output_dir / output_filename
        else:
            output_file = output_dir / (input_file.stem + ".md")

        if not markdown_text.strip():
            logger("      ⚠️ 警告：转换结果为空")
        output_file.write_text(markdown_text, encoding="utf-8")
        return True

    return False


def _process_single_file(
    file: Path,
    folder: Path,
    output_root: Path,
    temp_dir: Path,
    llm_client,
    llm_model: str | None,
    optimize: bool,
    analyze_attachments: bool,
) -> tuple[int, int, int]:
    """处理单个文件的辅助函数（用于并行处理）"""
    global processed_count

    # 计算相对路径
    rel_path = file.relative_to(folder)
    rel_dir = rel_path.parent

    # 立即记录开始处理，避免用户觉得卡住
    logger.info(f"🔄 正在处理: {file.name}")

    # 创建对应的输出目录结构
    current_output_dir = output_root / rel_dir
    assets_dir = current_output_dir / "assets"
    current_output_dir.mkdir(parents=True, exist_ok=True)

    working_file = file

    # 检查是否是需要转换的旧格式
    if file.suffix.lower() in OLD_FORMATS:
        # 使用全局锁保护 Office 转换，防止并发冲突
        with office_lock:
            working_file = convert_old_format(file, temp_dir, logger=logger.info)

        if working_file is None:
            with counter_lock:
                processed_count += 1
                curr_p = processed_count
            logger.warning(f"[{curr_p}/{total_count}] ⏭️  {file.name} 跳过 (转换失败)")
            return 0, 0, 1  # success, fail, skip

    # 在线程内部创建 MarkItDown 实例，确保线程安全
    llm_prompt = os.environ.get("LLM_PROMPT") or LLM_CONTENT_PROMPT

    file_ext = working_file.suffix.lower()
    if file_ext in IMAGE_FORMATS and llm_client:
        # 图片文件且有 LLM，使用带 LLM 的实例
        md = MarkItDown(llm_client=llm_client, llm_model=llm_model, llm_prompt=llm_prompt)
    else:
        # 其他情况使用普通实例
        md = MarkItDown()

    # 计算输出路径显示
    base_name = file.stem
    output_filename = f"{base_name}.md"
    output_file = current_output_dir / output_filename

    # 防重名逻辑
    counter = 1
    while output_file.exists():
        output_filename = f"{base_name}_{counter}.md"
        output_file = current_output_dir / output_filename
        counter += 1

    output_rel = rel_dir / output_filename if rel_dir != Path(".") else Path(output_filename)

    success = convert_to_markdown(
        working_file,
        current_output_dir,
        assets_dir,
        md,
        llm_client,
        llm_model,
        optimize,
        analyze_attachments,
        output_filename=output_filename,
        logger=logger.info,
    )

    # 处理完成后的记录
    with counter_lock:
        processed_count += 1
        curr_p = processed_count

    prefix = f"[{curr_p}/{total_count}]"
    if success:
        logger.info(f"{prefix} ✅ {file.name} -> {output_rel}")
        return 1, 0, 0
    else:
        logger.error(f"{prefix} ❌ {file.name} 失败")
        return 0, 1, 0


def process_folder(
    folder_path: str,
    optimize: bool = False,
    analyze_attachments: bool = False,
    analyze_image_files: bool = False,
    output_dir: str | None = None,
):
    """处理指定文件夹中的所有文件（包括子文件夹）

    Args:
        optimize: 使用 LLM 优化 Markdown 格式
        analyze_attachments: 使用 LLM 分析文档中提取出的附件图片
        analyze_image_files: 使用 LLM 分析图片文件
        output_dir: 指定输出目录（可选）
    """
    global total_count
    folder = Path(folder_path).resolve()

    if not folder.exists():
        logger.error(f"❌ 文件夹不存在: {folder_path}")
        return

    if not folder.is_dir():
        logger.error(f"❌ 路径不是文件夹: {folder_path}")
        return

    # 确定输出根目录
    env_output = os.environ.get("OUTPUT_DIR")

    if output_dir:
        output_root = Path(output_dir).resolve()
    elif env_output:
        output_root = Path(env_output).resolve()
    else:
        output_root = folder.parent / "output"

    temp_dir = output_root / "_temp"

    output_root.mkdir(parents=True, exist_ok=True)

    # 递归查找支持的文件
    all_files = []
    for f in folder.rglob("*"):
        if f.is_file() and f.suffix.lower() in SUPPORTED_FORMATS:
            try:
                f.relative_to(output_root)
                continue  # 文件在 output 目录中，跳过
            except ValueError:
                pass
            all_files.append(f)

    if not all_files:
        logger.warning(f"⚠️ 在 {folder} 中没有找到可转换的文件")
        logger.warning(f"   支持格式: {', '.join(sorted(SUPPORTED_FORMATS))}")
        return

    total_count = len(all_files)
    subdirs = {f.parent for f in all_files if f.parent != folder}
    image_files = [f for f in all_files if f.suffix.lower() in IMAGE_FORMATS]

    logger.info(f"📂 找到 {total_count} 个文件")
    if subdirs:
        logger.info(f"   包含 {len(subdirs)} 个子文件夹")
    logger.info(f"📁 输出: {output_root}")

    # 初始化 LLM（如果配置了且需要）
    llm_client = None
    llm_model = None
    should_init_llm = optimize or analyze_attachments or analyze_image_files

    if should_init_llm:
        llm_client, llm_model = create_llm_client()

    # 显示启用的功能
    if optimize:
        if llm_client:
            logger.info("✨ 已启用格式优化")
        else:
            logger.warning("⚠️ 格式优化需要配置 LLM，已跳过")
            optimize = False

    if analyze_attachments:
        if llm_client:
            logger.info("🔍 已启用附件图片分析")
        else:
            logger.warning("⚠️ 附件图片分析需要配置 LLM，已跳过")
            analyze_attachments = False

    if analyze_image_files:
        if llm_client:
            logger.info("🖼️ 已启用图片文件分析")
        else:
            logger.warning("⚠️ 图片文件分析需要配置 LLM，已跳过")
            analyze_image_files = False

    # 提示未启用的功能
    if image_files and not analyze_image_files:
        logger.info(
            f"💡 检测到 {len(image_files)} 个图片文件（使用 --analyze-image-files 启用 LLM 描述）"
        )

    logger.info("-" * 50)

    success, fail, skip = 0, 0, 0
    max_workers = int(os.environ.get("MAX_WORKERS", "10"))
    logger.info(f"🚀 启动 {max_workers} 个线程进行并行处理...")

    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_file = {
            executor.submit(
                _process_single_file,
                file,
                folder,
                output_root,
                temp_dir,
                llm_client,
                llm_model,
                optimize,
                analyze_attachments,
            ): file
            for file in sorted(all_files)
        }

        for future in concurrent.futures.as_completed(future_to_file):
            try:
                s, f, sk = future.result()
                success += s
                fail += f
                skip += sk
            except Exception as e:
                file_path = future_to_file[future]
                logger.error(f"❌ 处理文件异常 {file_path.name}: {e}")
                fail += 1

    # 清理
    if temp_dir.exists():
        shutil.rmtree(temp_dir)

    # 清理空的 assets 目录（递归）
    for assets in output_root.rglob("assets"):
        if assets.is_dir() and not any(assets.iterdir()):
            assets.rmdir()

    # 统计
    logger.info("=" * 50)
    logger.info(f"✨ 完成! 成功: {success}")
    if fail:
        logger.info(f"   失败: {fail}")
    if skip:
        logger.info(f"   跳过: {skip} (需要 LibreOffice 或转换失败)")
    logger.info(f"📁 输出位置: {output_root}")


def process_file(
    file_path: str,
    output_dir: str | None = None,
    optimize: bool | None = False,
    analyze_attachments: bool | None = False,
    analyze_image_files: bool | None = None,
):
    """处理单个文件"""
    global total_count
    input_file = Path(file_path).resolve()
    total_count = 1

    if not input_file.exists():
        logger.error(f"❌ 文件不存在: {file_path}")
        return

    if not input_file.is_file():
        logger.error(f"❌ 路径不是文件: {file_path}")
        return

    if input_file.suffix.lower() not in SUPPORTED_FORMATS:
        logger.error(f"❌ 不支持的文件格式: {input_file.suffix}")
        return

    # 自动判断
    is_image_file = input_file.suffix.lower() in IMAGE_FORMATS
    if analyze_image_files is None:
        analyze_image_files = is_image_file

    # 确定输出目录
    env_output = os.environ.get("OUTPUT_DIR")
    if output_dir:
        out_dir = Path(output_dir).resolve()
    elif env_output:
        out_dir = Path(env_output).resolve()
    else:
        out_dir = input_file.parent

    assets_dir = out_dir / "assets"
    temp_dir = out_dir / "_temp"

    out_dir.mkdir(parents=True, exist_ok=True)

    logger.info(f"📄 处理文件: {input_file.name}")
    logger.info(f"📁 输出目录: {out_dir}")

    # 初始化 LLM
    llm_client = None
    llm_model = None
    should_init_llm = optimize or analyze_attachments or analyze_image_files

    if should_init_llm:
        llm_client, llm_model = create_llm_client()

    # 显示启用的功能
    if optimize:
        if llm_client:
            logger.info("✨ 已启用格式优化")
        else:
            logger.warning("⚠️ 格式优化需要配置 LLM，已跳过")
            optimize = False

    if analyze_attachments:
        if llm_client:
            logger.info("🔍 已启用附件图片分析")
        else:
            logger.warning("⚠️ 附件图片分析需要配置 LLM，已跳过")
            analyze_attachments = False

    if analyze_image_files:
        if llm_client:
            logger.info("🖼️ 图片文件将使用 LLM 生成描述")
        else:
            logger.warning("⚠️ 图片文件分析需要配置 LLM，已跳过")
            analyze_image_files = False

    logger.info("-" * 50)

    # 获取自定义提示词
    llm_prompt = os.environ.get("LLM_PROMPT") or LLM_CONTENT_PROMPT

    # 根据文件类型选择 MarkItDown 实例
    file_ext = input_file.suffix.lower()
    use_llm_for_markitdown = analyze_image_files and llm_client and file_ext in IMAGE_FORMATS

    if use_llm_for_markitdown:
        md = MarkItDown(llm_client=llm_client, llm_model=llm_model, llm_prompt=llm_prompt)
    else:
        md = MarkItDown()

    # 处理旧格式
    working_file = input_file
    if input_file.suffix.lower() in OLD_FORMATS:
        working_file = convert_old_format(input_file, temp_dir, logger=logger.info)
        if working_file is None:
            return

    # 转换
    success = convert_to_markdown(
        working_file,
        out_dir,
        assets_dir,
        md,
        llm_client,
        llm_model,
        optimize,
        analyze_attachments,
        logger=logger.info,
    )

    if success:
        logger.info(f"✅ 转换成功: {out_dir / (input_file.stem + '.md')}")
    else:
        logger.error("❌ 转换失败")

    # 清理
    if temp_dir.exists():
        shutil.rmtree(temp_dir)
    if assets_dir.exists() and not any(assets_dir.iterdir()):
        assets_dir.rmdir()

    # 清理
    if temp_dir.exists():
        shutil.rmtree(temp_dir)
    if assets_dir.exists() and not any(assets_dir.iterdir()):
        assets_dir.rmdir()


def main():
    # 解决 Windows 下控制台编码问题（io 已在文件顶部导入）
    if sys.stdout.encoding != "utf-8":
        with contextlib.suppress(Exception):
            sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

    print("""
+===================================================+
|      Markdown 批量转换工具                        |
|      文档 → Markdown（含图片提取）                |
+===================================================+
    """)

    # 解析命令行参数
    args = sys.argv[1:]
    optimize = False
    analyze_attachments = False
    analyze_image_files = False
    analyze_image_files_explicit = False  # 是否显式指定了参数
    target_path = None
    output_dir = None

    # 检查 --optimize 或 -o 参数
    if "--optimize" in args:
        optimize = True
        args.remove("--optimize")
    if "-o" in args:
        optimize = True
        args.remove("-o")

    # 检查 --analyze-attachments 参数
    if "--analyze-attachments" in args:
        analyze_attachments = True
        args.remove("--analyze-attachments")

    # 检查 --analyze-image-files 参数
    if "--analyze-image-files" in args:
        analyze_image_files = True
        analyze_image_files_explicit = True
        args.remove("--analyze-image-files")

    # 检查 --output 或 -out 参数
    if "--output" in args:
        idx = args.index("--output")
        if idx + 1 < len(args):
            output_dir = args[idx + 1]
            args.pop(idx + 1)
            args.pop(idx)
    if "-out" in args:
        idx = args.index("-out")
        if idx + 1 < len(args):
            output_dir = args[idx + 1]
            args.pop(idx + 1)
            args.pop(idx)

    # 检查 --help 或 -h 参数
    if "--help" in args or "-h" in args:
        print('用法: uv run convert_to_markdown.py [选项] "文件或文件夹路径"')
        print()
        print("选项:")
        print("  -o, --optimize            使用 LLM 优化 Markdown 格式")
        print("  --analyze-attachments     使用 LLM 分析文档中提取的附件图片")
        print("  --analyze-image-files     使用 LLM 分析独立的图片文件")
        print("  -out, --output <dir>      指定输出目录")
        print("  -h, --help                显示此帮助信息")
        print()
        print("说明:")
        print("  默认不使用 LLM 处理图片")
        print("  直接指定图片文件时，--analyze-image-files 默认启用")
        print("  长文本（超过阈值）跳过所有 LLM 处理（分析与优化），直接返回原始内容")
        print("  所有选项可独立使用，也可组合使用")
        print()
        print("环境变量（用于 LLM 功能）:")
        print("  OPENAI_API_KEY      OpenAI API 密钥")
        print("  OPENAI_BASE_URL     自定义接口地址（如 OpenRouter）")
        print("  OPENAI_MODEL        模型名称")
        print("  GOOGLE_API_KEY      Google Gemini API 密钥")
        print("  OLLAMA_MODEL        本地 Ollama 模型名称")
        print("  LLM_PROMPT          自定义图片描述提示词")
        print("  LLM_CHUNK_SIZE      超长文本阈值（默认 10000 字符）")
        print("  MAX_WORKERS         文件并行处理线程数（默认 10）")
        print("  MAX_IMG_WORKERS     单文件内图片并行分析线程数（默认 5）")
        print("  OUTPUT_DIR          默认输出目录")
        print()
        print("示例:")
        print("  # 转换文件夹")
        print('  uv run convert_to_markdown.py "./documents"')
        print()
        print("  # 指定输出目录")
        print('  uv run convert_to_markdown.py -out "./output" "./documents"')
        print()
        print("  # 启用格式优化")
        print('  uv run convert_to_markdown.py -o "./documents"')
        return

    # 获取路径
    if args:
        target_path = args[0]
    else:
        print("请输入要转换的文件或文件夹路径（可拖入）:")
        target_path = input("> ").strip().strip('"').strip("'")

        # 交互模式下询问选项
        if target_path:
            print("\n是否启用 LLM 格式优化？(y/N):")
            optimize_input = input("> ").strip().lower()
            optimize = optimize_input in ("y", "yes", "是")

            print("\n是否启用附件图片分析？(y/N):")
            attach_input = input("> ").strip().lower()
            analyze_attachments = attach_input in ("y", "yes", "是")

            print("\n是否启用图片文件/PPTX 图片分析？(y/N):")
            img_input = input("> ").strip().lower()
            analyze_image_files = img_input in ("y", "yes", "是")
            analyze_image_files_explicit = True

    if not target_path:
        print("\n❌ 未提供路径")
        print('\n用法: uv run convert_to_markdown.py [选项] "文件或文件夹路径"')
        return

    # 判断是文件还是文件夹
    target = Path(target_path)
    if target.is_file():
        # 对于单个文件，如果未显式指定 analyze_image_files，则传 None（自动判断）
        img_files_arg = analyze_image_files if analyze_image_files_explicit else None
        process_file(
            target_path,
            output_dir=output_dir,
            optimize=optimize,
            analyze_attachments=analyze_attachments,
            analyze_image_files=img_files_arg,
        )
    elif target.is_dir():
        process_folder(
            target_path,
            output_dir=output_dir,
            optimize=optimize,
            analyze_attachments=analyze_attachments,
            analyze_image_files=analyze_image_files,
        )
    else:
        print(f"❌ 路径不存在: {target_path}")


if __name__ == "__main__":
    main()
