"""MarkItDown converter integration."""

import base64
import re
import zipfile
from pathlib import Path

import anyio

from markit.converters.base import BaseConverter, ConversionResult, ExtractedImage
from markit.exceptions import ConversionError
from markit.utils.logging import get_logger

log = get_logger(__name__)

# AutoShape placeholder message
AUTOSHAPE_PLACEHOLDER = (
    "\n\n> **Note:** This document contains AutoShape/VML graphics that could not be "
    "converted to Markdown. Please refer to the original source file for these elements.\n"
)

# Regex pattern for data URI: data:image/png;base64,<base64_data>
DATA_URI_PATTERN = re.compile(r"!\[([^\]]*)\]\((data:image/([^;]+);base64,([A-Za-z0-9+/=]+))\)")

# Regex pattern for markdown image syntax: ![alt](path)
IMAGE_REF_PATTERN = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")


def _normalize_image_spacing(markdown: str) -> str:
    """Normalize spacing around images.

    Ensures:
    - One blank line before first image (if preceded by text)
    - One blank line after last image (if followed by text)
    - At most one blank line between consecutive images
    """
    lines = markdown.split("\n")
    result = []
    last_content_type = "blank"  # "blank", "image", "text"

    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        is_image = stripped.startswith("![") and "](" in stripped
        is_blank = stripped == ""

        if is_image:
            if last_content_type == "image":
                # Consecutive images: keep as-is (no forced blank removal)
                result.append(line)
            elif last_content_type == "blank":
                # Already have blank line before image
                result.append(line)
            else:
                # Text before image: add blank line
                result.append("")
                result.append(line)
            last_content_type = "image"
        elif is_blank:
            # Limit to one blank line between images
            j = i + 1
            while j < len(lines) and lines[j].strip() == "":
                j += 1

            if j < len(lines):
                next_stripped = lines[j].strip()
                next_is_image = next_stripped.startswith("![") and "](" in next_stripped

                if last_content_type == "image" and next_is_image:
                    # Between consecutive images: keep at most one blank
                    result.append(line)
                    i = j - 1  # Skip extra blanks
                    last_content_type = "blank"
                elif last_content_type == "image":
                    # Keep one blank after image (before text)
                    result.append(line)
                    last_content_type = "blank"
                else:
                    result.append(line)
                    last_content_type = "blank"
            else:
                # End of file
                if last_content_type != "image":
                    result.append(line)
                last_content_type = "blank"
        else:
            # Regular text
            if last_content_type == "image":
                # Add blank line after image before text
                result.append("")
            result.append(line)
            last_content_type = "text"

        i += 1

    return "\n".join(result)


def _extract_pptx_footer_patterns(file_path: Path) -> set[str]:
    """Extract footer/header text patterns from PPTX file.

    Analyzes the PPTX structure to identify placeholder text that represents
    footers, headers, dates, and slide numbers.

    Args:
        file_path: Path to the PPTX file

    Returns:
        Set of text patterns that are likely footers/headers
    """
    if file_path.suffix.lower() != ".pptx":
        return set()

    patterns: set[str] = set()

    try:
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER

        # Placeholder types to filter
        FOOTER_PLACEHOLDER_TYPES = {
            PP_PLACEHOLDER.FOOTER,  # 15
            PP_PLACEHOLDER.DATE,  # 16
            PP_PLACEHOLDER.SLIDE_NUMBER,  # 13
            PP_PLACEHOLDER.HEADER,  # 12 (usually only in notes/handouts)
        }

        prs = Presentation(str(file_path))

        # Check slide masters for footer patterns
        for master in prs.slide_masters:
            for shape in master.shapes:
                if shape.is_placeholder:
                    try:
                        ph_type = shape.placeholder_format.type
                        if ph_type in FOOTER_PLACEHOLDER_TYPES:
                            if shape.has_text_frame:
                                text = shape.text.strip()  # type: ignore[attr-defined]
                                if text:
                                    patterns.add(text)
                    except Exception:
                        pass

        # Check slide layouts
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                for shape in layout.shapes:
                    if shape.is_placeholder:
                        try:
                            ph_type = shape.placeholder_format.type
                            if ph_type in FOOTER_PLACEHOLDER_TYPES:
                                if shape.has_text_frame:
                                    text = shape.text.strip()  # type: ignore[attr-defined]
                                    if text:
                                        patterns.add(text)
                        except Exception:
                            pass

        # Check actual slides for footer placeholders
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.is_placeholder:
                    try:
                        ph_type = shape.placeholder_format.type
                        if ph_type in FOOTER_PLACEHOLDER_TYPES:
                            if shape.has_text_frame:
                                text = shape.text.strip()  # type: ignore[attr-defined]
                                if text:
                                    patterns.add(text)
                    except Exception:
                        pass

    except ImportError:
        log.debug("python-pptx not available for footer detection")
    except Exception as e:
        log.debug("Could not extract PPTX footer patterns", file=str(file_path), error=str(e))

    return patterns


def _filter_pptx_footers(markdown: str, footer_patterns: set[str]) -> str:
    """Filter out footer/header text from PPTX markdown output.

    Args:
        markdown: Markdown content from PPTX conversion
        footer_patterns: Set of text patterns to filter out

    Returns:
        Filtered markdown content
    """
    if not footer_patterns:
        return markdown

    lines = markdown.split("\n")
    filtered_lines = []

    for line in lines:
        stripped = line.strip()
        # Skip lines that exactly match footer patterns
        if stripped in footer_patterns:
            continue
        # Also check for numeric-only lines (likely slide numbers)
        if stripped.isdigit():
            continue
        filtered_lines.append(line)

    return "\n".join(filtered_lines)


def _check_for_autoshapes(file_path: Path) -> bool:
    """Check if a document contains AutoShapes/VML graphics.

    Inspects DOCX/PPTX files for VML or DrawingML AutoShape elements
    that may not be fully converted to Markdown.

    Args:
        file_path: Path to the document file

    Returns:
        True if AutoShapes are detected, False otherwise
    """
    suffix = file_path.suffix.lower()
    if suffix not in (".docx", ".pptx"):
        return False

    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            for name in zf.namelist():
                # Check document.xml for DOCX or slide*.xml for PPTX
                if suffix == ".docx" and name == "word/document.xml":
                    content = zf.read(name).decode("utf-8", errors="ignore")
                    # VML shape detection
                    if "<w:pict>" in content or "<v:shape" in content:
                        return True
                    # DrawingML inline/anchor with shape (not picture)
                    if "<wps:wsp" in content:
                        return True
                elif (
                    suffix == ".pptx"
                    and name.startswith("ppt/slides/slide")
                    and name.endswith(".xml")
                ):
                    content = zf.read(name).decode("utf-8", errors="ignore")
                    # VML shape detection
                    if "<v:shape" in content:
                        return True
    except (zipfile.BadZipFile, OSError) as e:
        log.debug("Could not check for AutoShapes", file=str(file_path), error=str(e))

    return False


def _sanitize_filename(filename: str) -> str:
    """Sanitize filename by replacing special characters.

    Replaces characters that may cause issues in markdown links or file systems.
    """
    # Characters to replace with underscore
    replacements = {
        " ": "_",
        ":": "_",
        "：": "_",  # Full-width colon
        "/": "_",
        "\\": "_",
        "?": "_",
        "*": "_",
        '"': "_",
        "<": "_",
        ">": "_",
        "|": "_",
    }
    result = filename
    for char, replacement in replacements.items():
        result = result.replace(char, replacement)
    # Collapse multiple underscores
    while "__" in result:
        result = result.replace("__", "_")
    return result


class MarkItDownConverter(BaseConverter):
    """Converter using Microsoft's MarkItDown library.

    MarkItDown is the primary conversion engine that supports:
    - Text files (.txt)
    - Word documents (.docx)
    - PowerPoint presentations (.pptx)
    - Excel spreadsheets (.xlsx)
    - CSV files (.csv)
    - PDF documents (.pdf)
    - HTML files (.html, .htm)
    - Images (.png, .jpg, .jpeg) - with LLM support

    For legacy formats (.doc, .ppt, .xls), pre-processing with
    MS Office or LibreOffice is required first.
    """

    name = "markitdown"
    supported_extensions = {
        ".txt",
        ".docx",
        ".pptx",
        ".xlsx",
        ".csv",
        ".pdf",
        ".html",
        ".htm",
        ".png",
        ".jpg",
        ".jpeg",
        ".rtf",
        ".epub",
        ".zip",
        ".wav",
        ".mp3",
    }

    def __init__(self, enable_plugins: bool = True, keep_data_uris: bool = True) -> None:
        """Initialize MarkItDown converter.

        Args:
            enable_plugins: Whether to enable MarkItDown plugins
            keep_data_uris: Whether to keep images as data URIs (then extract them)
        """
        self.enable_plugins = enable_plugins
        self.keep_data_uris = keep_data_uris
        self._markitdown = None

    def _get_markitdown(self):
        """Lazy initialization of MarkItDown instance."""
        if self._markitdown is None:
            from markitdown import MarkItDown

            self._markitdown = MarkItDown(enable_plugins=self.enable_plugins)
        return self._markitdown

    async def convert(self, file_path: Path) -> ConversionResult:
        """Convert a document file to Markdown using MarkItDown.

        Args:
            file_path: Path to the document file

        Returns:
            ConversionResult with markdown content
        """
        if not await self.validate(file_path):
            raise ConversionError(
                file_path,
                f"File not supported or doesn't exist: {file_path.suffix}",
            )

        log.info("Converting with MarkItDown", file=str(file_path))

        try:
            # Run MarkItDown conversion in thread pool (it's synchronous)
            result = await anyio.to_thread.run_sync(self._convert_sync, file_path)  # type: ignore[attr-defined]
            return result
        except Exception as e:
            log.error(
                "MarkItDown conversion failed",
                file=str(file_path),
                error=str(e),
            )
            raise ConversionError(file_path, str(e), cause=e) from e

    def _convert_sync(self, file_path: Path) -> ConversionResult:
        """Synchronous conversion using MarkItDown."""
        md = self._get_markitdown()

        # Convert with keep_data_uris=True to get embedded images
        result = md.convert(str(file_path), keep_data_uris=self.keep_data_uris)

        # Extract markdown content
        # MarkItDown uses 'markdown' attribute (text_content is deprecated)
        markdown_content = getattr(result, "markdown", None)
        if markdown_content is None:
            # Fallback to text_content for older versions
            markdown_content = getattr(result, "text_content", "")

        # Get optional title
        title = getattr(result, "title", None)

        # Build metadata
        metadata = {
            "converter": self.name,
            "source_file": str(file_path),
            "file_type": file_path.suffix,
        }
        if title:
            metadata["title"] = title

        # Extract images from data URIs and local references
        images, markdown_content = self._extract_images(markdown_content, file_path)

        # Normalize spacing around images
        markdown_content = _normalize_image_spacing(markdown_content)

        # Post-process markdown for Excel files (remove NaN)
        if file_path.suffix.lower() in (".xlsx", ".xls", ".csv"):
            markdown_content = self._clean_excel_markdown(markdown_content)

        # Post-process PPTX: filter out footer/header placeholders
        if file_path.suffix.lower() == ".pptx":
            footer_patterns = _extract_pptx_footer_patterns(file_path)
            if footer_patterns:
                markdown_content = _filter_pptx_footers(markdown_content, footer_patterns)
                log.info(
                    "Filtered PPTX footer patterns",
                    file=str(file_path),
                    patterns=list(footer_patterns)[:5],  # Log first 5 patterns
                )

        # Check for AutoShapes and add placeholder if found
        warnings = []
        if _check_for_autoshapes(file_path):
            warnings.append("Document contains AutoShape/VML graphics that could not be converted")
            markdown_content += AUTOSHAPE_PLACEHOLDER
            log.info("AutoShape detected", file=str(file_path))

        return ConversionResult(
            markdown=markdown_content,
            images=images,
            metadata=metadata,
            success=True,
            warnings=warnings,
        )

    def _extract_images(self, markdown: str, source_file: Path) -> tuple[list[ExtractedImage], str]:
        """Extract images from markdown content.

        Extracts:
        1. Data URIs (base64 embedded images) - replaces with file references
        2. Local file references

        Returns:
            Tuple of (extracted images list, updated markdown with file references)
        """
        images: list[ExtractedImage] = []
        file_stem = source_file.stem

        # First, extract and replace data URIs
        def replace_data_uri(match: re.Match) -> str:
            nonlocal images
            alt_text = match.group(1)
            # full_uri = match.group(2)
            image_format = match.group(3)
            base64_data = match.group(4)

            # Decode base64 data
            try:
                image_data = base64.b64decode(base64_data)
            except Exception as e:
                log.warning("Failed to decode base64 image", error=str(e))
                return match.group(0)  # Return original

            # Normalize format
            if image_format == "jpg":
                image_format = "jpeg"
            elif image_format == "x-wmf":
                image_format = "wmf"
            elif image_format == "x-emf":
                image_format = "emf"

            # Generate filename with sanitized stem
            img_index = len(images) + 1
            safe_stem = _sanitize_filename(file_stem)
            filename = f"{safe_stem}_{img_index:03d}.{image_format}"

            images.append(
                ExtractedImage(
                    data=image_data,
                    format=image_format,
                    filename=filename,
                    source_document=source_file,
                    position=img_index,
                    original_path=None,
                )
            )

            # Return markdown with file reference
            return f"![{alt_text}](assets/{filename})"

        # Replace all data URIs
        markdown = DATA_URI_PATTERN.sub(replace_data_uri, markdown)

        # Then extract local file references
        for match in IMAGE_REF_PATTERN.finditer(markdown):
            alt_text = match.group(1)
            image_path = match.group(2)

            # Skip already processed (assets/) and external URLs
            if image_path.startswith(("data:", "http://", "https://", "assets/")):
                continue

            # Local path reference
            resolved_path = source_file.parent / image_path
            if resolved_path.exists():
                try:
                    with open(resolved_path, "rb") as f:
                        data = f.read()

                    # Determine format from extension
                    fmt = resolved_path.suffix.lower().lstrip(".")
                    if fmt == "jpg":
                        fmt = "jpeg"

                    img_index = len(images) + 1
                    safe_stem = _sanitize_filename(file_stem)
                    filename = f"{safe_stem}_{img_index:03d}.{fmt}"

                    images.append(
                        ExtractedImage(
                            data=data,
                            format=fmt,
                            filename=filename,
                            source_document=source_file,
                            position=img_index,
                            original_path=image_path,
                        )
                    )

                    # Update markdown reference
                    old_ref = f"![{alt_text}]({image_path})"
                    new_ref = f"![{alt_text}](assets/{filename})"
                    markdown = markdown.replace(old_ref, new_ref)

                except OSError:
                    log.warning(
                        "Could not read image file",
                        image_path=str(resolved_path),
                    )

        return images, markdown

    def _clean_excel_markdown(self, markdown: str) -> str:
        """Clean up Excel-generated markdown.

        - Replace NaN with empty string in table cells
        - Clean up table formatting
        """
        # Replace NaN in table cells with empty string
        # Pattern: | NaN | or | NaN at end of line
        # We need to be careful not to match "NaN" within actual content

        # Split by lines and process each line
        lines = markdown.split("\n")
        cleaned_lines = []

        for line in lines:
            if "|" in line:
                # Process table row - split by |, clean NaN, rejoin
                parts = line.split("|")
                cleaned_parts = []
                for part in parts:
                    # Check if this cell is exactly "NaN" (with optional whitespace)
                    stripped = part.strip()
                    if stripped == "NaN":
                        cleaned_parts.append(" ")
                    else:
                        cleaned_parts.append(part)
                line = "|".join(cleaned_parts)
            cleaned_lines.append(line)

        return "\n".join(cleaned_lines)


class MarkItDownPDFConverter(MarkItDownConverter):
    """MarkItDown converter specifically for PDF files.

    Uses MarkItDown's built-in PDF support.
    """

    name = "markitdown_pdf"
    supported_extensions = {".pdf"}


def get_markitdown_converter(enable_plugins: bool = True) -> MarkItDownConverter:
    """Factory function to create a MarkItDown converter.

    Args:
        enable_plugins: Whether to enable MarkItDown plugins

    Returns:
        Configured MarkItDownConverter instance
    """
    return MarkItDownConverter(enable_plugins=enable_plugins)
